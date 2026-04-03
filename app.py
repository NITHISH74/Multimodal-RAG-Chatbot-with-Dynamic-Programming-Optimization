"""
╔══════════════════════════════════════════════════════════════════════╗
║  Multi-Model RAG Chatbot — Streamlit Web Application               ║
║  Supports: Gemini & Cohere Embedding Models                        ║
║  Features: File Upload, Chat History, Token Tracking, DP Knapsack  ║
╚══════════════════════════════════════════════════════════════════════╝
"""

import os
import json
import time
import base64
import numpy as np
import faiss
import streamlit as st
from PIL import Image
from dotenv import load_dotenv
from google import genai
import cohere
from datetime import datetime

# ── Document parsers ──
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation

# ══════════════════════════════════════════════════════════════════════
#  CONFIGURATION & CONSTANTS
# ══════════════════════════════════════════════════════════════════════

load_dotenv()

FAISS_DIR = "faiss_data"
DATA_DIR = "data"
CHAT_HISTORY_FILE = os.path.join(FAISS_DIR, "chat_history.json")
GEMINI_EMBED_MODEL = "gemini-embedding-2-preview"
COHERE_EMBED_MODEL = "embed-v4.0"
GENERATION_MODEL = "gemini-2.5-flash"
MAX_CONTEXT_CHARS = 5000

os.makedirs(FAISS_DIR, exist_ok=True)
os.makedirs(DATA_DIR, exist_ok=True)

# ══════════════════════════════════════════════════════════════════════
#  PAGE CONFIGURATION & CUSTOM CSS
# ══════════════════════════════════════════════════════════════════════

st.set_page_config(
    page_title="Multi-Model RAG Chatbot",
    page_icon="🧠",
    layout="wide",
    initial_sidebar_state="expanded",
)

CUSTOM_CSS = """
<style>
    @import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700;800&display=swap');

    /* ── Global ── */
    .stApp {
        font-family: 'Inter', sans-serif;
    }
    
    /* ── Header ── */
    .main-header {
        background: linear-gradient(135deg, #0f0c29 0%, #302b63 50%, #24243e 100%);
        padding: 1.8rem 2rem;
        border-radius: 16px;
        margin-bottom: 1.5rem;
        box-shadow: 0 8px 32px rgba(48, 43, 99, 0.4);
        border: 1px solid rgba(255, 255, 255, 0.08);
    }
    .main-header h1 {
        color: #ffffff;
        font-weight: 800;
        font-size: 1.8rem;
        margin: 0;
        letter-spacing: -0.5px;
    }
    .main-header p {
        color: rgba(255,255,255,0.65);
        font-size: 0.9rem;
        margin: 0.3rem 0 0 0;
    }
    
    /* ── Metric Cards ── */
    .metric-row {
        display: flex;
        gap: 12px;
        margin-bottom: 1rem;
    }
    .metric-card {
        background: linear-gradient(135deg, rgba(255,255,255,0.05), rgba(255,255,255,0.02));
        backdrop-filter: blur(10px);
        border: 1px solid rgba(255,255,255,0.1);
        border-radius: 12px;
        padding: 1rem;
        flex: 1;
        text-align: center;
        transition: transform 0.2s ease, box-shadow 0.2s ease;
    }
    .metric-card:hover {
        transform: translateY(-2px);
        box-shadow: 0 8px 24px rgba(0,0,0,0.2);
    }
    .metric-value {
        font-size: 1.6rem;
        font-weight: 700;
        background: linear-gradient(135deg, #667eea, #764ba2);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        margin: 0;
    }
    .metric-label {
        font-size: 0.75rem;
        color: rgba(255,255,255,0.5);
        text-transform: uppercase;
        letter-spacing: 1px;
        margin: 0.3rem 0 0 0;
    }
    
    /* ── Chat Messages ── */
    .chat-user {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        color: white;
        padding: 1rem 1.3rem;
        border-radius: 18px 18px 4px 18px;
        margin: 0.6rem 0;
        max-width: 85%;
        margin-left: auto;
        box-shadow: 0 4px 16px rgba(102, 126, 234, 0.3);
        font-size: 0.92rem;
        line-height: 1.5;
    }
    .chat-bot {
        background: linear-gradient(135deg, rgba(255,255,255,0.08), rgba(255,255,255,0.03));
        border: 1px solid rgba(255,255,255,0.1);
        color: #e0e0e0;
        padding: 1rem 1.3rem;
        border-radius: 18px 18px 18px 4px;
        margin: 0.6rem 0;
        max-width: 85%;
        box-shadow: 0 4px 16px rgba(0,0,0,0.15);
        font-size: 0.92rem;
        line-height: 1.6;
    }
    
    /* ── Sidebar ── */
    [data-testid="stSidebar"] {
        background: linear-gradient(180deg, #0f0c29 0%, #1a1a2e 100%);
    }
    [data-testid="stSidebar"] .stMarkdown p {
        color: rgba(255,255,255,0.7);
    }
    
    .sidebar-section {
        background: rgba(255,255,255,0.04);
        border: 1px solid rgba(255,255,255,0.08);
        border-radius: 12px;
        padding: 1rem;
        margin: 0.8rem 0;
    }
    .sidebar-section h3 {
        font-size: 0.8rem;
        text-transform: uppercase;
        letter-spacing: 1.5px;
        color: rgba(255,255,255,0.4);
        margin: 0 0 0.6rem 0;
    }
    
    /* ── Status badges ── */
    .status-badge {
        display: inline-block;
        padding: 0.2rem 0.7rem;
        border-radius: 20px;
        font-size: 0.72rem;
        font-weight: 600;
        letter-spacing: 0.5px;
    }
    .status-online {
        background: rgba(46, 213, 115, 0.15);
        color: #2ed573;
        border: 1px solid rgba(46, 213, 115, 0.3);
    }
    .status-gemini {
        background: rgba(66, 133, 244, 0.15);
        color: #4285f4;
        border: 1px solid rgba(66, 133, 244, 0.3);
    }
    .status-cohere {
        background: rgba(168, 85, 247, 0.15);
        color: #a855f7;
        border: 1px solid rgba(168, 85, 247, 0.3);
    }
    
    /* ── File upload area ── */
    .upload-zone {
        border: 2px dashed rgba(102, 126, 234, 0.3);
        border-radius: 12px;
        padding: 1.2rem;
        text-align: center;
        background: rgba(102, 126, 234, 0.05);
        margin: 0.5rem 0;
        transition: border-color 0.3s ease;
    }
    .upload-zone:hover {
        border-color: rgba(102, 126, 234, 0.6);
    }
    
    /* ── Token bar ── */
    .token-bar-container {
        background: rgba(255,255,255,0.05);
        border-radius: 8px;
        padding: 0.6rem 0.8rem;
        margin: 0.5rem 0;
    }
    .token-bar {
        height: 6px;
        border-radius: 3px;
        background: linear-gradient(90deg, #667eea, #764ba2, #f093fb);
        transition: width 0.5s ease;
    }
    
    /* ── History item ── */
    .history-item {
        background: rgba(255,255,255,0.03);
        border: 1px solid rgba(255,255,255,0.06);
        border-radius: 8px;
        padding: 0.6rem 0.8rem;
        margin: 0.4rem 0;
        cursor: pointer;
        transition: background 0.2s ease;
    }
    .history-item:hover {
        background: rgba(255,255,255,0.08);
    }
    .history-item .time {
        font-size: 0.65rem;
        color: rgba(255,255,255,0.3);
    }
    .history-item .preview {
        font-size: 0.8rem;
        color: rgba(255,255,255,0.6);
        white-space: nowrap;
        overflow: hidden;
        text-overflow: ellipsis;
    }
    
    /* ── Animations ── */
    @keyframes fadeIn {
        from { opacity: 0; transform: translateY(8px); }
        to { opacity: 1; transform: translateY(0); }
    }
    .animate-in {
        animation: fadeIn 0.4s ease-out;
    }
    
    /* ── Scrollbar ── */
    ::-webkit-scrollbar { width: 6px; }
    ::-webkit-scrollbar-track { background: transparent; }
    ::-webkit-scrollbar-thumb { background: rgba(255,255,255,0.15); border-radius: 3px; }
    ::-webkit-scrollbar-thumb:hover { background: rgba(255,255,255,0.25); }
    
    /* ── Hide default streamlit elements ── */
    #MainMenu { visibility: hidden; }
    footer { visibility: hidden; }
    .stDeployButton { display: none; }
</style>
"""

st.markdown(CUSTOM_CSS, unsafe_allow_html=True)


# ══════════════════════════════════════════════════════════════════════
#  STATE INITIALIZATION
# ══════════════════════════════════════════════════════════════════════

def init_state():
    """Initialize Streamlit session state variables."""
    defaults = {
        "messages": [],
        "total_input_tokens": 0,
        "total_output_tokens": 0,
        "total_queries": 0,
        "embedding_model": "Gemini",
        "gemini_client": None,
        "cohere_client": None,
        "faiss_index": None,
        "faiss_metadata": [],
        "indexed_count": 0,
        "chat_history_all": [],
        "current_session_id": datetime.now().strftime("%Y%m%d_%H%M%S"),
    }
    for key, val in defaults.items():
        if key not in st.session_state:
            st.session_state[key] = val

init_state()


# ══════════════════════════════════════════════════════════════════════
#  CLIENT INITIALIZATION
# ══════════════════════════════════════════════════════════════════════

def get_gemini_client():
    if st.session_state.gemini_client is None:
        api_key = os.getenv("gemini_api_key")
        if api_key:
            st.session_state.gemini_client = genai.Client(api_key=api_key)
    return st.session_state.gemini_client


def get_cohere_client():
    if st.session_state.cohere_client is None:
        api_key = os.getenv("cohere_api_key")
        if api_key:
            st.session_state.cohere_client = cohere.ClientV2(api_key=api_key)
    return st.session_state.cohere_client


# ══════════════════════════════════════════════════════════════════════
#  FAISS INDEX MANAGEMENT
# ══════════════════════════════════════════════════════════════════════

def get_index_path(model_name):
    """Each embedding model gets its own FAISS index."""
    prefix = "gemini" if model_name == "Gemini" else "cohere"
    return (
        os.path.join(FAISS_DIR, f"{prefix}_index.bin"),
        os.path.join(FAISS_DIR, f"{prefix}_metadata.json"),
    )


def load_faiss_index(model_name):
    """Load or create FAISS index for the selected embedding model."""
    index_file, meta_file = get_index_path(model_name)
    default_dim = 3072 if model_name == "Gemini" else 1536

    try:
        if os.path.exists(index_file) and os.path.exists(meta_file):
            index = faiss.read_index(index_file)
            with open(meta_file, "r", encoding="utf-8") as f:
                metadata = json.load(f)
            return index, metadata
    except Exception:
        pass

    index = faiss.IndexFlatL2(default_dim)
    return index, []


def save_faiss_index(model_name, index, metadata):
    """Persist FAISS index and metadata to disk."""
    index_file, meta_file = get_index_path(model_name)
    faiss.write_index(index, index_file)
    with open(meta_file, "w", encoding="utf-8") as f:
        json.dump(metadata, f, indent=2, ensure_ascii=False)


# ══════════════════════════════════════════════════════════════════════
#  FILE PARSERS
# ══════════════════════════════════════════════════════════════════════

def extract_text_from_pdf(file_bytes):
    """Extract text from a PDF file."""
    import io
    reader = PdfReader(io.BytesIO(file_bytes))
    text = ""
    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            text += page_text + "\n"
    return text.strip()


def extract_text_from_docx(file_bytes):
    """Extract text from a DOCX file."""
    import io
    doc = DocxDocument(io.BytesIO(file_bytes))
    return "\n".join(para.text for para in doc.paragraphs if para.text.strip())


def extract_text_from_pptx(file_bytes):
    """Extract text from a PPTX file."""
    import io
    prs = Presentation(io.BytesIO(file_bytes))
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
    return "\n".join(texts)


def parse_uploaded_file(uploaded_file):
    """
    Parse an uploaded file and return (content_type, content, filename).
    content_type: 'text' or 'image'
    content: extracted text string or PIL Image
    """
    name = uploaded_file.name.lower()
    file_bytes = uploaded_file.read()

    if name.endswith(".pdf"):
        text = extract_text_from_pdf(file_bytes)
        return "text", text, uploaded_file.name
    elif name.endswith(".docx"):
        text = extract_text_from_docx(file_bytes)
        return "text", text, uploaded_file.name
    elif name.endswith(".pptx"):
        text = extract_text_from_pptx(file_bytes)
        return "text", text, uploaded_file.name
    elif name.endswith((".txt", ".md")):
        text = file_bytes.decode("utf-8", errors="ignore")
        return "text", text, uploaded_file.name
    elif name.endswith((".png", ".jpg", ".jpeg", ".webp", ".gif")):
        import io
        img = Image.open(io.BytesIO(file_bytes))
        # Also save to data/ for persistence
        save_path = os.path.join(DATA_DIR, uploaded_file.name)
        with open(save_path, "wb") as f:
            f.write(file_bytes)
        return "image", img, uploaded_file.name
    else:
        return None, None, uploaded_file.name


# ══════════════════════════════════════════════════════════════════════
#  EMBEDDING FUNCTIONS
# ══════════════════════════════════════════════════════════════════════

def embed_text_gemini(text):
    """Embed text using Gemini. Returns vector."""
    client = get_gemini_client()
    response = client.models.embed_content(model=GEMINI_EMBED_MODEL, contents=text)
    return response.embeddings[0].values


def embed_image_gemini(img):
    """Embed PIL Image using Gemini. Returns vector."""
    client = get_gemini_client()
    response = client.models.embed_content(model=GEMINI_EMBED_MODEL, contents=img)
    return response.embeddings[0].values


def embed_text_cohere(text, input_type="search_query"):
    """Embed text using Cohere. Returns vector."""
    client = get_cohere_client()
    response = client.embed(
        model=COHERE_EMBED_MODEL,
        texts=[text],
        input_type=input_type,
        embedding_types=["float"],
    )
    return list(response.embeddings.float_[0])


def embed_image_cohere(image_path):
    """Embed image using Cohere via base64. Returns vector."""
    client = get_cohere_client()
    ext = os.path.splitext(image_path)[1].lower().replace(".", "")
    mime_map = {"jpg": "jpeg", "jpeg": "jpeg", "png": "png", "webp": "webp", "gif": "gif"}
    mime = mime_map.get(ext, "jpeg")
    with open(image_path, "rb") as f:
        encoded = base64.b64encode(f.read()).decode("utf-8")
    data_url = f"data:image/{mime};base64,{encoded}"
    response = client.embed(
        model=COHERE_EMBED_MODEL,
        images=[data_url],
        input_type="image",
        embedding_types=["float"],
    )
    return list(response.embeddings.float_[0])


def embed_content(content, content_type, model_name, filename=None):
    """
    Unified embedding function.
    Returns the embedding vector based on model and content type.
    """
    if model_name == "Gemini":
        if content_type == "text":
            return embed_text_gemini(content)
        elif content_type == "image":
            return embed_image_gemini(content)
    else:  # Cohere
        if content_type == "text":
            return embed_text_cohere(content, input_type="search_document")
        elif content_type == "image":
            img_path = os.path.join(DATA_DIR, filename) if filename else None
            if img_path and os.path.exists(img_path):
                return embed_image_cohere(img_path)
            else:
                raise ValueError(f"Image file not found for Cohere embedding: {filename}")
    return None


def embed_query(query_text, model_name):
    """Embed a query for retrieval."""
    if model_name == "Gemini":
        return embed_text_gemini(query_text)
    else:
        return embed_text_cohere(query_text, input_type="search_query")


# ══════════════════════════════════════════════════════════════════════
#  DP 0/1 KNAPSACK — CONTEXT OPTIMIZATION
# ══════════════════════════════════════════════════════════════════════

def optimize_context_knapsack(candidates, max_chars):
    """Classic DP 0/1 Knapsack to select optimal context within token budget."""
    n = len(candidates)
    if n == 0:
        return []

    SCALE = 10
    W = max_chars // SCALE
    weights = [max(1, item["weight"] // SCALE) for item in candidates]
    values = [item["value"] for item in candidates]

    dp = [[0.0] * (W + 1) for _ in range(n + 1)]
    for i in range(1, n + 1):
        for w in range(1, W + 1):
            if weights[i - 1] <= w:
                dp[i][w] = max(dp[i - 1][w], dp[i - 1][w - weights[i - 1]] + values[i - 1])
            else:
                dp[i][w] = dp[i - 1][w]

    chosen = []
    w = W
    for i in range(n, 0, -1):
        if dp[i][w] != dp[i - 1][w]:
            chosen.append(candidates[i - 1])
            w -= weights[i - 1]

    return chosen


# ══════════════════════════════════════════════════════════════════════
#  CHAT HISTORY PERSISTENCE
# ══════════════════════════════════════════════════════════════════════

def save_chat_history():
    """Save chat history to disk."""
    history_entry = {
        "session_id": st.session_state.current_session_id,
        "timestamp": datetime.now().isoformat(),
        "embedding_model": st.session_state.embedding_model,
        "messages": st.session_state.messages,
        "total_input_tokens": st.session_state.total_input_tokens,
        "total_output_tokens": st.session_state.total_output_tokens,
        "total_queries": st.session_state.total_queries,
    }

    # Load existing history
    all_history = load_all_chat_history()

    # Update or append current session
    updated = False
    for i, h in enumerate(all_history):
        if h["session_id"] == st.session_state.current_session_id:
            all_history[i] = history_entry
            updated = True
            break
    if not updated:
        all_history.append(history_entry)

    # Keep only last 50 sessions
    all_history = all_history[-50:]

    with open(CHAT_HISTORY_FILE, "w", encoding="utf-8") as f:
        json.dump(all_history, f, indent=2, ensure_ascii=False)


def load_all_chat_history():
    """Load all chat history sessions."""
    if os.path.exists(CHAT_HISTORY_FILE):
        try:
            with open(CHAT_HISTORY_FILE, "r", encoding="utf-8") as f:
                return json.load(f)
        except Exception:
            return []
    return []


def load_session(session_id):
    """Load a specific chat session."""
    history = load_all_chat_history()
    for h in history:
        if h["session_id"] == session_id:
            return h
    return None


# ══════════════════════════════════════════════════════════════════════
#  RAG PIPELINE
# ══════════════════════════════════════════════════════════════════════

def run_rag_pipeline(user_query, model_name):
    """
    Full RAG pipeline:
    1. Embed query
    2. Search FAISS
    3. DP Knapsack optimization
    4. Generate response with Gemini
    Returns (response_text, metadata_dict)
    """
    index, metadata = load_faiss_index(model_name)
    pipeline_meta = {
        "retrieval_time": 0,
        "generation_time": 0,
        "context_chunks": 0,
        "input_tokens": 0,
        "output_tokens": 0,
    }

    candidates = []
    t0 = time.time()

    if index.ntotal > 0:
        try:
            query_vec = embed_query(user_query, model_name)
            query_np = np.array([query_vec]).astype("float32")

            k = min(10, index.ntotal)
            distances, indices = index.search(query_np, k)

            for idx, dist in zip(indices[0], distances[0]):
                if idx != -1:
                    matched = metadata[idx]
                    val = 1.0 / (dist + 0.001)
                    candidates.append({
                        "file": matched["file"],
                        "type": matched["type"],
                        "content": matched["content"],
                        "value": float(val),
                        "weight": matched.get("weight", len(matched.get("content", ""))),
                        "dist": float(dist),
                    })
        except Exception as e:
            st.warning(f"Retrieval error: {e}")

    pipeline_meta["retrieval_time"] = time.time() - t0

    # DP Knapsack optimization
    chosen = optimize_context_knapsack(candidates, MAX_CONTEXT_CHARS)
    pipeline_meta["context_chunks"] = len(chosen)

    # Build augmented prompt
    t1 = time.time()
    contexts = []
    image_parts = []

    for item in chosen:
        if item["type"] == "text":
            contexts.append(f"Context from {item['file']} (relevance {item['value']:.2f}):\n{item['content']}\n")
        elif item["type"] == "image":
            img_path = os.path.join(DATA_DIR, item["file"])
            if os.path.exists(img_path):
                try:
                    img = Image.open(img_path)
                    image_parts.append(img)
                    contexts.append(f"Context from {item['file']}: (Image attached)\n")
                except Exception:
                    pass

    prompt = user_query
    if contexts:
        context_str = "\n".join(contexts)
        prompt = (
            f"Below is optimized context retrieved from my local knowledge base (DP 0/1 Knapsack filtered):\n"
            f"----------\n{context_str}\n----------\n"
            f"Based on the context and your knowledge, answer the user's query:\n\n"
            f"{user_query}"
        )

    # Estimate input tokens (rough: 1 token ≈ 4 chars)
    pipeline_meta["input_tokens"] = len(prompt) // 4

    # Generate with Gemini
    try:
        client = get_gemini_client()
        message_parts = [prompt]
        message_parts.extend(image_parts)

        response = client.models.generate_content(
            model=GENERATION_MODEL,
            contents=message_parts,
        )
        response_text = response.text
        
        # Try to get actual token usage from response
        if hasattr(response, 'usage_metadata') and response.usage_metadata:
            um = response.usage_metadata
            if hasattr(um, 'prompt_token_count') and um.prompt_token_count:
                pipeline_meta["input_tokens"] = um.prompt_token_count
            if hasattr(um, 'candidates_token_count') and um.candidates_token_count:
                pipeline_meta["output_tokens"] = um.candidates_token_count
        else:
            pipeline_meta["output_tokens"] = len(response_text) // 4

    except Exception as e:
        response_text = f"⚠️ Generation error: {e}"
        pipeline_meta["output_tokens"] = 0

    pipeline_meta["generation_time"] = time.time() - t1

    return response_text, pipeline_meta


# ══════════════════════════════════════════════════════════════════════
#  INDEXING LOGIC
# ══════════════════════════════════════════════════════════════════════

def index_uploaded_files(uploaded_files, model_name, progress_bar, status_text):
    """Index uploaded files into the FAISS store."""
    index, metadata = load_faiss_index(model_name)
    indexed_files = {item["file"] for item in metadata}

    new_vectors = []
    new_metadata = []
    total = len(uploaded_files)

    for i, uploaded_file in enumerate(uploaded_files):
        fname = uploaded_file.name
        progress_bar.progress((i + 1) / total, text=f"Processing {fname}...")
        status_text.markdown(f"🔄 Embedding `{fname}` with **{model_name}**...")

        if fname in indexed_files:
            status_text.markdown(f"⏩ `{fname}` already indexed, skipping.")
            continue

        try:
            content_type, content, filename = parse_uploaded_file(uploaded_file)
            if content_type is None:
                status_text.warning(f"❌ Unsupported file type: {fname}")
                continue

            if content_type == "text" and not content.strip():
                status_text.warning(f"❌ No text extracted from: {fname}")
                continue

            # Save text files to data/ for persistence
            if content_type == "text":
                save_path = os.path.join(DATA_DIR, fname)
                with open(save_path, "w", encoding="utf-8") as f:
                    f.write(content)

            vector = embed_content(content, content_type, model_name, filename=fname)

            if vector:
                # Resize index if dimension mismatch on first insert
                if index.ntotal == 0 and len(vector) != index.d:
                    index = faiss.IndexFlatL2(len(vector))

                new_vectors.append(vector)
                meta_entry = {
                    "file": fname,
                    "type": content_type,
                    "content": content if content_type == "text" else "Visual content",
                    "weight": len(content) if content_type == "text" else 500,
                }
                new_metadata.append(meta_entry)
                status_text.markdown(f"✅ `{fname}` embedded successfully!")

        except Exception as e:
            status_text.error(f"❌ Error processing {fname}: {e}")

    if new_vectors:
        embeddings_np = np.array(new_vectors).astype("float32")
        index.add(embeddings_np)
        metadata.extend(new_metadata)
        save_faiss_index(model_name, index, metadata)

    return len(new_vectors)


def index_existing_data(model_name, progress_bar, status_text):
    """Index files already in the data/ folder."""
    index, metadata = load_faiss_index(model_name)
    indexed_files = {item["file"] for item in metadata}

    files = [f for f in os.listdir(DATA_DIR) if f not in indexed_files]
    if not files:
        return 0

    new_vectors = []
    new_metadata = []

    for i, fname in enumerate(files):
        filepath = os.path.join(DATA_DIR, fname)
        progress_bar.progress((i + 1) / len(files), text=f"Processing {fname}...")
        status_text.markdown(f"🔄 Embedding `{fname}` with **{model_name}**...")

        try:
            if fname.lower().endswith((".txt", ".md")):
                with open(filepath, "r", encoding="utf-8") as f:
                    content = f.read()
                vector = embed_content(content, "text", model_name)
                meta = {"file": fname, "type": "text", "content": content, "weight": len(content)}

            elif fname.lower().endswith((".png", ".jpg", ".jpeg", ".webp")):
                img = Image.open(filepath)
                vector = embed_content(img, "image", model_name, filename=fname)
                meta = {"file": fname, "type": "image", "content": "Visual content", "weight": 500}

            elif fname.lower().endswith(".pdf"):
                with open(filepath, "rb") as f:
                    content = extract_text_from_pdf(f.read())
                if not content.strip():
                    continue
                vector = embed_content(content, "text", model_name)
                meta = {"file": fname, "type": "text", "content": content, "weight": len(content)}

            elif fname.lower().endswith(".docx"):
                with open(filepath, "rb") as f:
                    content = extract_text_from_docx(f.read())
                if not content.strip():
                    continue
                vector = embed_content(content, "text", model_name)
                meta = {"file": fname, "type": "text", "content": content, "weight": len(content)}

            elif fname.lower().endswith(".pptx"):
                with open(filepath, "rb") as f:
                    content = extract_text_from_pptx(f.read())
                if not content.strip():
                    continue
                vector = embed_content(content, "text", model_name)
                meta = {"file": fname, "type": "text", "content": content, "weight": len(content)}

            else:
                continue

            if vector:
                if index.ntotal == 0 and len(vector) != index.d:
                    index = faiss.IndexFlatL2(len(vector))
                new_vectors.append(vector)
                new_metadata.append(meta)
                status_text.markdown(f"✅ `{fname}` embedded!")

        except Exception as e:
            status_text.error(f"❌ Error: {fname}: {e}")

    if new_vectors:
        embeddings_np = np.array(new_vectors).astype("float32")
        index.add(embeddings_np)
        metadata.extend(new_metadata)
        save_faiss_index(model_name, index, metadata)

    return len(new_vectors)


# ══════════════════════════════════════════════════════════════════════
#  UI LAYOUT
# ══════════════════════════════════════════════════════════════════════

# ── SIDEBAR ──────────────────────────────────────────────────────────

with st.sidebar:
    # Logo / Brand
    st.markdown("""
    <div style="text-align: center; padding: 1rem 0 0.5rem 0;">
        <span style="font-size: 2.5rem;">🧠</span>
        <h2 style="margin: 0.3rem 0 0 0; font-weight: 700; font-size: 1.2rem; 
                    background: linear-gradient(135deg, #667eea, #764ba2);
                    -webkit-background-clip: text; -webkit-text-fill-color: transparent;">
            Multi-Model RAG
        </h2>
        <p style="color: rgba(255,255,255,0.4); font-size: 0.7rem; margin: 0;">
            Powered by Gemini 2.5 Flash
        </p>
    </div>
    """, unsafe_allow_html=True)

    st.markdown("---")

    # ── Embedding Model Selector ──
    st.markdown("""
    <div class="sidebar-section">
        <h3>⚡ Embedding Engine</h3>
    </div>
    """, unsafe_allow_html=True)

    model_choice = st.radio(
        "Select Embedding Model",
        ["Gemini", "Cohere"],
        index=0 if st.session_state.embedding_model == "Gemini" else 1,
        key="model_radio",
        help="Switch between Gemini Embedding 2 Preview and Cohere Embed v4.0",
    )

    if model_choice != st.session_state.embedding_model:
        st.session_state.embedding_model = model_choice
        st.rerun()

    model_badge = "status-gemini" if model_choice == "Gemini" else "status-cohere"
    model_label = GEMINI_EMBED_MODEL if model_choice == "Gemini" else COHERE_EMBED_MODEL
    st.markdown(f"""
    <span class="status-badge {model_badge}">{model_label}</span>
    """, unsafe_allow_html=True)

    st.markdown("---")

    # ── File Upload ──
    st.markdown("""
    <div class="sidebar-section">
        <h3>📁 Upload Documents</h3>
    </div>
    """, unsafe_allow_html=True)

    uploaded_files = st.file_uploader(
        "Upload files to the knowledge base",
        type=["txt", "md", "pdf", "docx", "pptx", "png", "jpg", "jpeg", "webp"],
        accept_multiple_files=True,
        key="file_uploader",
        help="Supported: TXT, MD, PDF, DOCX, PPTX, PNG, JPG, JPEG, WebP",
    )

    col_idx1, col_idx2 = st.columns(2)
    with col_idx1:
        index_uploaded_btn = st.button("📥 Index Uploads", use_container_width=True, key="idx_uploads")
    with col_idx2:
        index_existing_btn = st.button("🔄 Index Data/", use_container_width=True, key="idx_existing")

    # Index uploaded files
    if index_uploaded_btn and uploaded_files:
        progress = st.progress(0, text="Starting...")
        status = st.empty()
        count = index_uploaded_files(uploaded_files, st.session_state.embedding_model, progress, status)
        if count > 0:
            st.success(f"✅ Indexed {count} new files with {st.session_state.embedding_model}!")
        else:
            st.info("No new files to index.")
        time.sleep(1)
        st.rerun()

    # Index existing data/ files
    if index_existing_btn:
        progress = st.progress(0, text="Starting...")
        status = st.empty()
        count = index_existing_data(st.session_state.embedding_model, progress, status)
        if count > 0:
            st.success(f"✅ Indexed {count} files from data/ with {st.session_state.embedding_model}!")
        else:
            st.info("No new files to index from data/.")
        time.sleep(1)
        st.rerun()

    # Show indexed stats
    idx, meta = load_faiss_index(st.session_state.embedding_model)
    st.markdown(f"""
    <div class="sidebar-section">
        <h3>📊 Index Stats</h3>
        <p style="margin: 0.2rem 0;">
            <span style="color: #667eea; font-weight: 600;">{idx.ntotal}</span> 
            documents indexed
        </p>
        <p style="margin: 0.2rem 0; font-size: 0.75rem; color: rgba(255,255,255,0.4);">
            Model: {st.session_state.embedding_model} | Dim: {idx.d}
        </p>
    </div>
    """, unsafe_allow_html=True)

    st.markdown("---")

    # ── Token Usage ──
    total_tokens = st.session_state.total_input_tokens + st.session_state.total_output_tokens
    st.markdown(f"""
    <div class="sidebar-section">
        <h3>🔢 Token Usage</h3>
        <div style="display: flex; justify-content: space-between; margin: 0.3rem 0;">
            <span style="color: rgba(255,255,255,0.5); font-size: 0.8rem;">Input Tokens</span>
            <span style="color: #667eea; font-weight: 600; font-size: 0.8rem;">{st.session_state.total_input_tokens:,}</span>
        </div>
        <div style="display: flex; justify-content: space-between; margin: 0.3rem 0;">
            <span style="color: rgba(255,255,255,0.5); font-size: 0.8rem;">Output Tokens</span>
            <span style="color: #764ba2; font-weight: 600; font-size: 0.8rem;">{st.session_state.total_output_tokens:,}</span>
        </div>
        <div style="display: flex; justify-content: space-between; margin: 0.3rem 0;">
            <span style="color: rgba(255,255,255,0.5); font-size: 0.8rem;">Total</span>
            <span style="color: #f093fb; font-weight: 700; font-size: 0.85rem;">{total_tokens:,}</span>
        </div>
        <div class="token-bar-container">
            <div class="token-bar" style="width: {min(100, total_tokens / 100)}%;"></div>
        </div>
        <p style="margin: 0.3rem 0 0 0; font-size: 0.7rem; color: rgba(255,255,255,0.3);">
            Queries this session: {st.session_state.total_queries}
        </p>
    </div>
    """, unsafe_allow_html=True)

    st.markdown("---")

    # ── Chat History ──
    st.markdown("""
    <div class="sidebar-section">
        <h3>🕐 Chat History</h3>
    </div>
    """, unsafe_allow_html=True)

    all_history = load_all_chat_history()
    if all_history:
        for h in reversed(all_history[-10:]):
            sid = h["session_id"]
            ts = h.get("timestamp", "")[:16].replace("T", " ")
            n_msgs = len(h.get("messages", []))
            model_used = h.get("embedding_model", "?")
            preview = ""
            if h.get("messages"):
                first_user = next((m["content"][:40] for m in h["messages"] if m["role"] == "user"), "")
                preview = first_user + "..." if first_user else ""

            is_current = sid == st.session_state.current_session_id
            border_color = "rgba(102, 126, 234, 0.4)" if is_current else "rgba(255,255,255,0.06)"

            if st.button(f"{'🟢 ' if is_current else ''}{ts} ({n_msgs} msgs, {model_used})", key=f"hist_{sid}", use_container_width=True):
                if not is_current:
                    session_data = load_session(sid)
                    if session_data:
                        st.session_state.messages = session_data.get("messages", [])
                        st.session_state.total_input_tokens = session_data.get("total_input_tokens", 0)
                        st.session_state.total_output_tokens = session_data.get("total_output_tokens", 0)
                        st.session_state.total_queries = session_data.get("total_queries", 0)
                        st.session_state.current_session_id = sid
                        st.rerun()
    else:
        st.markdown('<p style="color: rgba(255,255,255,0.3); font-size: 0.8rem;">No history yet.</p>', unsafe_allow_html=True)

    # ── New Chat / Clear ──
    st.markdown("---")
    col_c1, col_c2 = st.columns(2)
    with col_c1:
        if st.button("🆕 New Chat", use_container_width=True, key="new_chat_btn"):
            if st.session_state.messages:
                save_chat_history()
            st.session_state.messages = []
            st.session_state.total_input_tokens = 0
            st.session_state.total_output_tokens = 0
            st.session_state.total_queries = 0
            st.session_state.current_session_id = datetime.now().strftime("%Y%m%d_%H%M%S")
            st.rerun()
    with col_c2:
        if st.button("🗑️ Clear All", use_container_width=True, key="clear_all_btn"):
            if os.path.exists(CHAT_HISTORY_FILE):
                os.remove(CHAT_HISTORY_FILE)
            st.session_state.messages = []
            st.session_state.total_input_tokens = 0
            st.session_state.total_output_tokens = 0
            st.session_state.total_queries = 0
            st.rerun()


# ── MAIN CONTENT ──────────────────────────────────────────────────────

# Header
model_emoji = "🟢" if st.session_state.embedding_model == "Gemini" else "🟣"
st.markdown(f"""
<div class="main-header animate-in">
    <h1>🧠 Multi-Model RAG Chatbot</h1>
    <p>Gemini 2.5 Flash Generation • {model_emoji} {st.session_state.embedding_model} Embeddings • DP Knapsack Optimization • FAISS Vector Search</p>
</div>
""", unsafe_allow_html=True)

# Metrics row
idx, meta = load_faiss_index(st.session_state.embedding_model)
st.markdown(f"""
<div class="metric-row animate-in">
    <div class="metric-card">
        <p class="metric-value">{idx.ntotal}</p>
        <p class="metric-label">Documents Indexed</p>
    </div>
    <div class="metric-card">
        <p class="metric-value">{st.session_state.total_queries}</p>
        <p class="metric-label">Queries Made</p>
    </div>
    <div class="metric-card">
        <p class="metric-value">{st.session_state.total_input_tokens + st.session_state.total_output_tokens:,}</p>
        <p class="metric-label">Total Tokens</p>
    </div>
    <div class="metric-card">
        <p class="metric-value">{st.session_state.embedding_model}</p>
        <p class="metric-label">Embedding Model</p>
    </div>
</div>
""", unsafe_allow_html=True)


# ── Chat Display ──
chat_container = st.container()
with chat_container:
    for msg in st.session_state.messages:
        role = msg["role"]
        content = msg["content"]
        if role == "user":
            with st.chat_message("user", avatar="👤"):
                st.markdown(content)
        else:
            with st.chat_message("assistant", avatar="🧠"):
                st.markdown(content)
                # Show metadata if available
                if "meta" in msg:
                    m = msg["meta"]
                    cols = st.columns(4)
                    cols[0].caption(f"⏱️ Retrieval: {m.get('retrieval_time', 0):.2f}s")
                    cols[1].caption(f"⚡ Generation: {m.get('generation_time', 0):.2f}s")
                    cols[2].caption(f"📎 Chunks: {m.get('context_chunks', 0)}")
                    cols[3].caption(f"🔢 Tokens: {m.get('input_tokens', 0) + m.get('output_tokens', 0):,}")


# ── Chat Input ──
if user_input := st.chat_input("Ask anything about your documents...", key="chat_input"):
    # Add user message
    st.session_state.messages.append({"role": "user", "content": user_input})

    with st.chat_message("user", avatar="👤"):
        st.markdown(user_input)

    # Generate response
    with st.chat_message("assistant", avatar="🧠"):
        with st.spinner(f"🔍 Searching with {st.session_state.embedding_model} embeddings..."):
            response_text, pipeline_meta = run_rag_pipeline(user_input, st.session_state.embedding_model)

        st.markdown(response_text)

        # Show pipeline metadata
        cols = st.columns(4)
        cols[0].caption(f"⏱️ Retrieval: {pipeline_meta['retrieval_time']:.2f}s")
        cols[1].caption(f"⚡ Generation: {pipeline_meta['generation_time']:.2f}s")
        cols[2].caption(f"📎 Chunks: {pipeline_meta['context_chunks']}")
        cols[3].caption(f"🔢 Tokens: {pipeline_meta['input_tokens'] + pipeline_meta['output_tokens']:,}")

    # Update state
    st.session_state.messages.append({
        "role": "assistant",
        "content": response_text,
        "meta": pipeline_meta,
    })
    st.session_state.total_input_tokens += pipeline_meta["input_tokens"]
    st.session_state.total_output_tokens += pipeline_meta["output_tokens"]
    st.session_state.total_queries += 1

    # Auto-save
    save_chat_history()
    st.rerun()
